import os
import re
import uuid
import threading
import asyncio
import pygame
import edge_tts
import speech_recognition as sr
import webbrowser
import subprocess
import wikipedia
import difflib
import time
import random
import shutil
from docx import Document
from pptx import Presentation
from PIL import Image, ImageFilter

# ---------------- SETTINGS ----------------
VOICE = "en-AU-WilliamNeural"
AUTHORIZED_USER = "Vashu"
AUTH_RETRY = 3
# Updated QNA file path
QNA_FILE = "C:/Users/ACER/Desktop/pyhon/jarvis/data/brain_data/qna.txt"
os.makedirs(os.path.dirname(QNA_FILE), exist_ok=True)

# ---------------- TTS ----------------
pygame.mixer.init()
tts_lock = threading.Lock()

def remove_file(file_path):
    try:
        os.remove(file_path)
    except:
        pass

async def tts_save(text, filename):
    communicator = edge_tts.Communicate(text, VOICE)
    await communicator.save(filename)
    threading.Thread(target=play_audio, args=(filename,), daemon=True).start()

def play_audio(file_path):
    try:
        with tts_lock:
            sound = pygame.mixer.Sound(file_path)
            sound.play()
        while pygame.mixer.get_busy():
            pygame.time.delay(100)
        remove_file(file_path)
    except:
        pass

def speak(text):
    filename = f"voice_{uuid.uuid4().hex}.mp3"
    try:
        asyncio.run(tts_save(text, filename))
    except:
        print("TTS error")

# ---------------- LISTEN ----------------
def listen(timeout=6):
    r = sr.Recognizer()
    try:
        with sr.Microphone() as source:
            print("Listening...")
            r.adjust_for_ambient_noise(source, duration=0.5)
            audio = r.listen(source, phrase_time_limit=timeout)
            try:
                text = r.recognize_google(audio)
                print(f"You said: {text}")
                return text
            except:
                pass
    except:
        pass
    return ""

# ---------------- HELPERS ----------------
def ask_input(prompt, default=""):
    speak(prompt)
    ans = listen().strip()
    if not ans:
        print(prompt)
        ans = input("Type here (or press Enter to skip): ").strip()
    return ans if ans else default

def safe_int(s, default=0):
    try:
        return int(s)
    except:
        return default

# ---------------- SECURITY ----------------
def authenticate_user():
    for i in range(AUTH_RETRY):
        speak("Please say your name for authentication")
        name = listen(timeout=5).strip()
        if not name:
            name = input("Type your name: ").strip()
        similarity = difflib.SequenceMatcher(None, name.lower(), AUTHORIZED_USER.lower()).ratio()
        if similarity > 0.7:
            speak(f"Welcome back, {name}!")
            return True
        speak(f"Authentication failed {i+1}/{AUTH_RETRY}")
    speak("Unauthorized. Exiting.")
    return False

# ---------------- WISH & WELCOME ----------------
def make_wish():
    h = int(time.strftime("%H"))
    if h < 12:
        speak("Good Morning!")
    elif h < 18:
        speak("Good Afternoon!")
    else:
        speak("Good Evening!")

def welcome():
    greetings = ["Hello!", "Hi there!", "Welcome back!"]
    speak(random.choice(greetings))

# ---------------- QNA ----------------
def load_qna(file_path=QNA_FILE):
    qna_dict = {}
    if os.path.exists(file_path):
        with open(file_path, "r", encoding="utf-8") as f:
            for line in f:
                if ":" in line:
                    q, a = line.strip().split(":", 1)
                    qna_dict[q.lower().strip()] = a.strip()
    return qna_dict

def save_qna(question, answer, file_path=QNA_FILE):
    with open(file_path, "a", encoding="utf-8") as f:
        f.write(f"{question}: {answer}\n")

def answer_question(user_question, qna_dict):
    best_match = difflib.get_close_matches(user_question.lower(), qna_dict.keys(), n=1, cutoff=0.6)
    if best_match:
        return qna_dict[best_match[0]]
    return None

# ---------------- PATH HELPERS ----------------
USER_HOME = os.path.expanduser("~")
USER_FOLDERS = {
    "desktop": os.path.join(USER_HOME, "Desktop"),
    "documents": os.path.join(USER_HOME, "Documents"),
    "downloads": os.path.join(USER_HOME, "Downloads"),
    "pictures": os.path.join(USER_HOME, "Pictures"),
    "music": os.path.join(USER_HOME, "Music"),
    "videos": os.path.join(USER_HOME, "Videos")
}

def search_and_open(target):
    for folder in USER_FOLDERS.values():
        if os.path.exists(folder):
            for root, dirs, files in os.walk(folder):
                for item in dirs + files:
                    if target.lower() in item.lower():
                        path = os.path.join(root, item)
                        os.startfile(path)
                        speak(f"Opening {item}")
                        return True
    return False

# ---------------- WEB & SEARCH ----------------
KNOWN_SITES = {
    "youtube": "https://www.youtube.com",
    "google": "https://www.google.com",
    "gmail": "https://mail.google.com"
}

def handle_open(command):
    target = re.sub(r'(open|launch|start|go to)', '', command.lower()).strip()
    if target in ["my pc", "this pc", "computer"]:
        subprocess.Popen("explorer.exe shell:MyComputerFolder")
        return
    if target in USER_FOLDERS:
        os.startfile(USER_FOLDERS[target])
        return
    if target in KNOWN_SITES:
        webbrowser.open_new_tab(KNOWN_SITES[target])
        return
    if search_and_open(target):
        return
    exe = target.replace(" ", "") + ".exe"
    path = shutil.which(exe) or shutil.which(target)
    if path:
        subprocess.Popen(path)
        return
    webbrowser.open_new_tab(f"https://www.google.com/search?q={'+'.join(target.split())}")

def handle_play(cmd):
    song = re.sub(r'(play|music|songs?|video|by|of|on youtube)', '', cmd, flags=re.I).strip()
    if not song:
        speak("What to play?")
        return
    url = "https://www.youtube.com/results?search_query=" + "+".join(song.split())
    webbrowser.open_new_tab(url)
    speak(f"Searching {song} on YouTube.")

def system_command(cmd):
    try:
        if cmd == "shutdown":
            os.system("shutdown /s /t 5")
        elif cmd == "restart":
            os.system("shutdown /r /t 5")
        elif cmd == "lock":
            os.system("rundll32.exe user32.dll,LockWorkStation")
        elif cmd == "logout":
            os.system("shutdown /l")
    except:
        pass

# ---------------- SMART CV ----------------
def smart_cv(filename=None):
    try:
        name = ask_input("Please say or type your full name", "Candidate")
        email = ask_input("Please provide your email", "not_provided@example.com")
        phone = ask_input("Please provide your phone number", "0000000000")
        role = ask_input("Please provide the job role", "Desired Role")
        skills = ask_input("List your skills, comma separated", "No skills provided")
        summary = ask_input("Provide a short summary (or skip)", f"{name} applying for {role}, skilled in {skills}.")
        experience = ask_input("Describe your experience (or skip)", f"{name} has experience in {role}.")
        education = ask_input("Mention your education (or skip)", "")
        projects = ask_input("Mention any projects (or skip)", "")
        if not filename:
            filename = f"Smart_CV_{name.replace(' ', '_')}.docx"
        doc = Document()
        doc.add_heading("Curriculum Vitae", 0)
        doc.add_heading("Personal Info", 1)
        doc.add_paragraph(f"Name: {name}\nEmail: {email}\nPhone: {phone}")
        doc.add_heading("Summary", 1); doc.add_paragraph(summary)
        doc.add_heading("Experience", 1); doc.add_paragraph(experience)
        doc.add_heading("Skills", 1); doc.add_paragraph(skills)
        if education:
            doc.add_heading("Education", 1); doc.add_paragraph(education)
        if projects:
            doc.add_heading("Projects", 1); doc.add_paragraph(projects)
        doc.save(filename)
        speak(f"Smart CV saved as {filename}")
        print("CV saved at:", filename)
        try:
            os.startfile(filename)
        except:
            pass
    except Exception as e:
        speak("Failed to create CV")
        print("smart_cv error:", e)

# ---------------- PPT CREATOR ----------------
def generate_ppt_web(topic):
    try:
        page = wikipedia.page(topic, auto_suggest=True)
        content = page.content.split(". ")
        prs = Presentation()
        slides_count = min(8, max(3, len(content)//5))
        for i in range(slides_count):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"{topic} - Slide {i+1}"
            text_chunk = ". ".join(content[i*5:(i+1)*5])
            slide.placeholders[1].text = text_chunk
        filename = f"PPT_{topic.replace(' ', '_')}.pptx"
        prs.save(filename)
        speak(f"PPT on {topic} created as {filename}")
        print("PPT saved at:", filename)
        try:
            os.startfile(filename)
        except:
            pass
    except Exception as e:
        speak("Failed to generate PPT from web")
        print("generate_ppt_web error:", e)

# ---------------- MAIN ----------------
def wiki_search(query):
    try:
        q = re.sub(r'(who is|what is|tell me about|define|jarvis)', '', query.lower()).strip().title()
        if not q:
            return None
        return wikipedia.summary(q, sentences=2)
    except:
        return None

def jarvis():
    if not authenticate_user():
        return
    make_wish()
    welcome()
    qna_dict = load_qna()
    while True:
        user_input = listen()
        if not user_input:
            user_input = input("You (type command): ").strip()
        if not user_input:
            continue
        cmd = user_input.lower().strip()

        if cmd in ["exit", "quit", "stop", "goodbye"]:
            speak("Goodbye!")
            break
        if "smart cv" in cmd:
            smart_cv()
            continue
        if "create ppt" in cmd or "generate ppt" in cmd:
            topic = ask_input("Provide the topic for PPT", "Demo Topic")
            generate_ppt_web(topic)
            continue

        for sc in ["shutdown", "restart", "lock", "logout"]:
            if sc in cmd:
                system_command(sc)
                break

        if cmd.startswith("play"):
            handle_play(cmd)
            continue
        if cmd.startswith(("open", "launch", "start", "go to")):
            handle_open(cmd)
            continue

        response = answer_question(cmd, qna_dict)
        if response:
            speak(response)
        else:
            wiki = wiki_search(cmd)
            if wiki:
                speak(wiki)
                save_qna(cmd, wiki)
                qna_dict[cmd.lower()] = wiki
            else:
                webbrowser.open_new_tab(f"https://www.google.com/search?q={'+'.join(cmd.split())}")

if __name__ == "__main__":
    jarvis()
