import os
import re
import webbrowser
import subprocess
import wikipedia
import difflib
import shutil
from flask import Flask, render_template, request, jsonify, send_from_directory
from docx import Document
from pptx import Presentation

# ---------------- SETTINGS ----------------
AUTHORIZED_USER = "Vashu"
QNA_FILE = "C:/Users/ACER/Desktop/pyhon/jarvis/data/brain_data/qna.txt"
os.makedirs(os.path.dirname(QNA_FILE), exist_ok=True)

USER_HOME = os.path.expanduser("~")
USER_FOLDERS = {
    "desktop": os.path.join(USER_HOME, "Desktop"),
    "documents": os.path.join(USER_HOME, "Documents"),
    "downloads": os.path.join(USER_HOME, "Downloads"),
    "pictures": os.path.join(USER_HOME, "Pictures"),
    "music": os.path.join(USER_HOME, "Music"),
    "videos": os.path.join(USER_HOME, "Videos")
}

KNOWN_SITES = {
    "youtube": "https://www.youtube.com",
    "google": "https://www.google.com",
    "gmail": "https://mail.google.com"
}

GENERATED_FOLDER = os.path.join(os.getcwd(), "generated")
os.makedirs(GENERATED_FOLDER, exist_ok=True)

# ---------------- QNA ----------------
def load_qna(file_path=QNA_FILE):
    qna_dict = {}
    if os.path.exists(file_path):
        with open(file_path, "r", encoding="utf-8") as f:
            for line in f:
                if ":" in line:
                    q, a = line.strip().split(":", 1)
                    qna_dict[q.lower().strip()] = a.strip()
    return qna_dict

def save_qna(question, answer, file_path=QNA_FILE):
    with open(file_path, "a", encoding="utf-8") as f:
        f.write(f"{question}: {answer}\n")

def answer_question(user_question, qna_dict):
    best_match = None
    best_ratio = 0
    for q in qna_dict.keys():
        ratio = difflib.SequenceMatcher(None, user_question.lower(), q.lower()).ratio()
        if ratio > best_ratio:
            best_ratio = ratio
            best_match = q
    if best_match and best_ratio >= 0.75:
        return qna_dict[best_match]
    return None

# ---------------- PATH HELPERS ----------------
def open_path(path):
    if os.path.exists(path):
        if os.name == "nt":
            os.startfile(path)
        elif os.name == "posix":
            subprocess.Popen(["xdg-open", path])
        return f"Opening {os.path.basename(path)}"
    return None

def search_and_open(target):
    for folder in USER_FOLDERS.values():
        for root, dirs, files in os.walk(folder):
            for item in dirs + files:
                if target.lower() in item.lower():
                    return open_path(os.path.join(root, item))
    return None

# ---------------- COMMAND HANDLERS ----------------
def handle_open(command):
    target = re.sub(r'\b(open|launch|start|go to)\b', '', command.lower()).strip()

    # Open known websites
    for site, url in KNOWN_SITES.items():
        if site in target:
            webbrowser.open_new_tab(url)
            return f"Opening {site}"

    # Open folders
    for folder_name, folder_path in USER_FOLDERS.items():
        if folder_name in target:
            return open_path(folder_path)

    # Search local files/folders
    result = search_and_open(target)
    if result:
        return result

    # Try executables
    exe = target.replace(" ","") + ".exe"
    path = shutil.which(exe) or shutil.which(target)
    if path:
        subprocess.Popen(path)
        return f"Launching {target}"

    # Fallback: Google search
    webbrowser.open_new_tab(f"https://www.google.com/search?q={'+'.join(target.split())}")
    return f"Searching Google for {target}"

def handle_play(cmd):
    song = re.sub(r'(play|music|songs?|video|by|of|on youtube)', '', cmd, flags=re.I).strip()
    if not song:
        return "No song specified"
    url = "https://www.youtube.com/results?search_query=" + "+".join(song.split())
    webbrowser.open_new_tab(url)
    return f"Searching YouTube for {song}"

def system_command(cmd):
    try:
        if cmd == "shutdown":
            os.system("shutdown /s /t 5")
        elif cmd == "restart":
            os.system("shutdown /r /t 5")
        elif cmd == "lock":
            os.system("rundll32.exe user32.dll,LockWorkStation")
        elif cmd == "logout":
            os.system("shutdown /l")
        return f"Executed {cmd}"
    except:
        return "System command failed"

# ---------------- SMART CV ----------------
def smart_cv():
    filename = os.path.join(GENERATED_FOLDER, "Smart_CV.docx")
    try:
        doc = Document()
        doc.add_heading("Curriculum Vitae", 0)
        doc.add_heading("Personal Info", 1)
        doc.add_paragraph("Name: Candidate\nEmail: not_provided@example.com\nPhone: 0000000000")
        doc.add_heading("Summary", 1)
        doc.add_paragraph("Candidate applying for Desired Role, skilled in No skills provided.")
        doc.add_heading("Experience", 1)
        doc.add_paragraph("Candidate has experience in Desired Role.")
        doc.add_heading("Skills", 1)
        doc.add_paragraph("No skills provided")
        doc.save(filename)
        return f"Smart CV saved as {filename}"
    except Exception as e:
        return f"Failed to create Smart CV: {e}"

# ---------------- PPT CREATOR ----------------
def generate_ppt_web(topic):
    filename = os.path.join(GENERATED_FOLDER, f"PPT_{topic.replace(' ','_')}.pptx")
    try:
        page = wikipedia.page(topic, auto_suggest=True)
        content = page.content.split(". ")
        prs = Presentation()
        slides_count = min(8, max(3, len(content)//5))
        for i in range(slides_count):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"{topic} - Slide {i+1}"
            text_chunk = ". ".join(content[i*5:(i+1)*5])
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape.text_frame.text = text_chunk
                    break
        prs.save(filename)
        return f"PPT created as {filename}"
    except (wikipedia.DisambiguationError, wikipedia.PageError):
        return f"Topic '{topic}' not found on Wikipedia"
    except Exception as e:
        return f"Failed to generate PPT: {e}"

# ---------------- WIKIPEDIA ----------------
def wiki_search(query):
    try:
        cleaned = re.sub(r'(who is|what is|tell me about|define|jarvis|please|explain)', '', query.lower()).strip()
        if not cleaned:
            return None
        summary = wikipedia.summary(cleaned, sentences=2)
        return summary
    except (wikipedia.DisambiguationError, wikipedia.PageError):
        return None
    except Exception:
        return None

# ---------------- FLASK ----------------
app = Flask(__name__)

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/process', methods=['POST'])
def process():
    data = request.get_json()
    cmd = data.get('command', '').strip().lower()
    topic = data.get('topic', 'Demo Topic')

    if not cmd:
        return jsonify({'reply': 'Please say something.'})

    qna_dict = load_qna()

    # Smart CV
    if "smart cv" in cmd:
        result = smart_cv()
        return jsonify({'reply': result})

    # Generate PPT
    if "create ppt" in cmd or "generate ppt" in cmd:
        result = generate_ppt_web(topic)
        return jsonify({'reply': result})

    # System commands
    for sc in ["shutdown", "restart", "lock", "logout"]:
        if sc in cmd:
            result = system_command(sc)
            return jsonify({'reply': result})

    # Play music/video
    if cmd.startswith("play"):
        result = handle_play(cmd)
        return jsonify({'reply': result})

    # Open websites/folders/files
    if cmd.startswith(("open","launch","start","go to")):
        result = handle_open(cmd)
        return jsonify({'reply': result})

    # QnA / Wikipedia / Google fallback
    response = answer_question(cmd, qna_dict)
    if response:
        return jsonify({'reply': response})
    else:
        wiki = wiki_search(cmd)
        if wiki:
            save_qna(cmd, wiki)
            return jsonify({'reply': wiki})
        else:
            webbrowser.open_new_tab(f"https://www.google.com/search?q={'+'.join(cmd.split())}")
            return jsonify({'reply': f"Searching Google for {cmd}"})

# Download generated files
@app.route('/download/<filename>')
def download_file(filename):
    return send_from_directory(GENERATED_FOLDER, filename, as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True)
